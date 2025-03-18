from typing import Dict, List, Optional, Any, BinaryIO
from fastapi import HTTPException, status
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, and_, func, desc
from datetime import datetime, date, timedelta
import json
import io
import uuid
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.colors import grey, whitesmoke, RGBColor
from pptx import Presentation
from pptx.util import Inches

from models.db_models import TrainingSession, Attendance, Feedback, User
from models.report import (
    TrainingReportResponse,
    AttendanceReportResponse,
    FeedbackReportResponse,
    TrainingSessionSummary,
    AttendanceStats,
    AttendanceSummary,
    FeedbackSummary,
    SessionAttendance,
    AthleteAttendanceReport,
    TeamAttendanceReport,
    SessionFeedbackReport,
    FeedbackDetail,
    ReportFormat
)


class ReportService:
    """Service for generating reports"""
    
    def __init__(self, db_session: AsyncSession):
        self.db_session = db_session
    
    async def generate_training_report(
        self,
        start_date: datetime,
        end_date: datetime,
        user_id: str,
        user_role: str
    ) -> TrainingReportResponse:
        """Generate a training report for a date range"""
        # Build query based on date range
        query = select(TrainingSession).where(
            and_(
                TrainingSession.date >= start_date,
                TrainingSession.date <= end_date
            )
        )
        
        # Apply role-based filtering
        if user_role == "coach":
            query = query.where(TrainingSession.coach_id == user_id)
        
        # Order by date
        query = query.order_by(TrainingSession.date.desc())
        
        # Execute query
        result = await self.db_session.execute(query)
        sessions = result.scalars().all()
        
        # Get session summaries
        session_summaries = []
        total_attendees = 0
        
        for session in sessions:
            # Get attendance count for this session
            attendance_query = select(func.count()).where(
                Attendance.training_session_id == session.id
            )
            attendance_result = await self.db_session.execute(attendance_query)
            attendees_count = attendance_result.scalar()
            total_attendees += attendees_count
            
            # Get coach name
            coach_query = select(User).where(User.id == session.coach_id)
            coach_result = await self.db_session.execute(coach_query)
            coach = coach_result.scalars().first()
            
            # Get feedback for this session
            feedback_query = select(Feedback).where(
                Feedback.training_session_id == session.id
            )
            feedback_result = await self.db_session.execute(feedback_query)
            feedbacks = feedback_result.scalars().all()
            
            # Calculate average ratings
            feedback_count = len(feedbacks)
            training_quality_avg = 0
            expectations_avg = 0
            body_condition_avg = 0
            intensity_avg = 0
            
            if feedback_count > 0:
                training_quality_avg = sum(fb.training_quality for fb in feedbacks) / feedback_count
                expectations_avg = sum(fb.expectations for fb in feedbacks) / feedback_count
                body_condition_avg = sum(fb.body_condition for fb in feedbacks) / feedback_count
                intensity_avg = sum(fb.intensity for fb in feedbacks) / feedback_count
            
            # Create session summary
            session_summary = TrainingSessionSummary(
                id=session.id,
                date=session.date.strftime("%Y-%m-%d"),
                type=session.type,
                coach_name=coach.name if coach else "Unknown",
                duration_minutes=int((session.end_time - session.start_time).total_seconds() / 60),
                attendees_count=attendees_count,
                feedback={
                    "training_quality_avg": round(training_quality_avg, 1),
                    "expectations_avg": round(expectations_avg, 1),
                    "body_condition_avg": round(body_condition_avg, 1),
                    "intensity_avg": round(intensity_avg, 1)
                }
            )
            
            session_summaries.append(session_summary)
        
        # Create report summary
        total_sessions = len(sessions)
        avg_attendees_per_session = round(total_attendees / total_sessions, 1) if total_sessions > 0 else 0
        
        summary = {
            "total_sessions": total_sessions,
            "total_attendees": total_attendees,
            "avg_attendees_per_session": avg_attendees_per_session,
            "date_range": f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}"
        }
        
        # Create the full report response
        report = TrainingReportResponse(
            title="Training Sessions Report",
            generated_at=datetime.now(),
            date_range=f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
            summary=summary,
            data=session_summaries
        )
        
        return report
    
    async def generate_attendance_report(
        self,
        start_date: datetime,
        end_date: datetime,
        athlete_id: Optional[str] = None
    ) -> AttendanceReportResponse:
        """Generate an attendance report for a date range"""
        # Get sessions in the date range
        session_query = select(TrainingSession).where(
            and_(
                TrainingSession.date >= start_date,
                TrainingSession.date <= end_date
            )
        ).order_by(TrainingSession.date.desc())
        
        session_result = await self.db_session.execute(session_query)
        sessions = session_result.scalars().all()
        
        # Report title and summary depends on whether we're reporting for a single athlete or the whole team
        if athlete_id:
            # Individual athlete report
            return await self._generate_athlete_attendance_report(
                athlete_id=athlete_id,
                sessions=sessions,
                start_date=start_date,
                end_date=end_date
            )
        else:
            # Team report
            return await self._generate_team_attendance_report(
                sessions=sessions,
                start_date=start_date,
                end_date=end_date
            )
    
    async def _generate_athlete_attendance_report(
        self,
        athlete_id: str,
        sessions: List[TrainingSession],
        start_date: datetime,
        end_date: datetime
    ) -> AttendanceReportResponse:
        """Generate an attendance report for a specific athlete"""
        # Get athlete info
        athlete_query = select(User).where(User.id == athlete_id)
        athlete_result = await self.db_session.execute(athlete_query)
        athlete = athlete_result.scalars().first()
        
        if not athlete:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Athlete not found"
            )
        
        # Check attendance for each session
        session_attendances = []
        attended_count = 0
        
        for session in sessions:
            # Check if athlete attended this session
            attendance_query = select(Attendance).where(
                and_(
                    Attendance.training_session_id == session.id,
                    Attendance.athlete_id == athlete_id
                )
            )
            attendance_result = await self.db_session.execute(attendance_query)
            attendance = attendance_result.scalars().first()
            
            # Get coach name
            coach_query = select(User).where(User.id == session.coach_id)
            coach_result = await self.db_session.execute(coach_query)
            coach = coach_result.scalars().first()
            
            # Create session attendance record
            session_attendance = SessionAttendance(
                id=session.id,
                date=session.date.strftime("%Y-%m-%d"),
                type=session.type,
                coach_name=coach.name if coach else "Unknown",
                attended=attendance is not None,
                check_in_time=attendance.check_in_time.strftime("%H:%M:%S") if attendance else None
            )
            
            session_attendances.append(session_attendance)
            
            if attendance:
                attended_count += 1
        
        # Calculate attendance rate
        total_count = len(sessions)
        attendance_rate = (attended_count / total_count * 100) if total_count > 0 else 0
        
        # Create athlete attendance report
        athlete_report = AthleteAttendanceReport(
            athlete_id=athlete_id,
            athlete_name=athlete.name,
            attended_count=attended_count,
            total_count=total_count,
            attendance_rate=round(attendance_rate, 1),
            sessions=session_attendances
        )
        
        # Create report summary
        summary = {
            "athlete_name": athlete.name,
            "attended_count": attended_count,
            "total_count": total_count,
            "attendance_rate": round(attendance_rate, 1)
        }
        
        # Create the full report response
        report = AttendanceReportResponse(
            title=f"Attendance Report for {athlete.name}",
            generated_at=datetime.now(),
            date_range=f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
            summary=summary,
            data=[athlete_report]
        )
        
        return report
    
    async def _generate_team_attendance_report(
        self,
        sessions: List[TrainingSession],
        start_date: datetime,
        end_date: datetime
    ) -> AttendanceReportResponse:
        """Generate an attendance report for the whole team"""
        # Get all athletes (only those with role="athlete")
        athlete_query = select(User).where(User.role == "athlete")
        athlete_result = await self.db_session.execute(athlete_query)
        athletes = athlete_result.scalars().all()
        
        # Check attendance for each athlete across all sessions
        attendance_summaries = []
        total_attendance_rate = 0
        
        for athlete in athletes:
            # Count attended sessions for this athlete
            attended_count = 0
            
            for session in sessions:
                # Check if athlete attended this session
                attendance_query = select(Attendance).where(
                    and_(
                        Attendance.training_session_id == session.id,
                        Attendance.athlete_id == athlete.id
                    )
                )
                attendance_result = await self.db_session.execute(attendance_query)
                attendance = attendance_result.scalars().first()
                
                if attendance:
                    attended_count += 1
            
            # Calculate attendance rate
            total_count = len(sessions)
            attendance_rate = (attended_count / total_count * 100) if total_count > 0 else 0
            total_attendance_rate += attendance_rate
            
            # Create attendance summary
            attendance_summary = AttendanceSummary(
                athlete_id=athlete.id,
                athlete_name=athlete.name,
                sessions_attended=attended_count,
                sessions_missed=total_count - attended_count,
                attendance_rate=round(attendance_rate, 1)
            )
            
            attendance_summaries.append(attendance_summary)
        
        # Sort summaries by attendance rate (descending)
        attendance_summaries.sort(key=lambda x: x.attendance_rate, reverse=True)
        
        # Calculate team average
        avg_attendance_rate = (total_attendance_rate / len(athletes)) if athletes else 0
        
        # Create team attendance report
        team_report = TeamAttendanceReport(
            total_athletes=len(athletes),
            avg_attendance_rate=round(avg_attendance_rate, 1),
            athletes=attendance_summaries
        )
        
        # Create report summary
        summary = {
            "total_athletes": len(athletes),
            "total_sessions": len(sessions),
            "avg_attendance_rate": round(avg_attendance_rate, 1)
        }
        
        # Create the full report response
        report = AttendanceReportResponse(
            title="Team Attendance Report",
            generated_at=datetime.now(),
            date_range=f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
            summary=summary,
            data=attendance_summaries
        )
        
        return report
    
    async def generate_feedback_report(
        self,
        start_date: datetime,
        end_date: datetime,
        session_id: Optional[str] = None
    ) -> FeedbackReportResponse:
        """Generate a feedback report for a date range or specific session"""
        if session_id:
            # Report for a specific session
            return await self._generate_session_feedback_report(
                session_id=session_id,
                start_date=start_date,
                end_date=end_date
            )
        else:
            # Summary report for all sessions in date range
            return await self._generate_feedback_summary_report(
                start_date=start_date,
                end_date=end_date
            )
    
    async def _generate_session_feedback_report(
        self,
        session_id: str,
        start_date: datetime,
        end_date: datetime
    ) -> FeedbackReportResponse:
        """Generate a detailed feedback report for a specific session"""
        # Get session info
        session_query = select(TrainingSession).where(TrainingSession.id == session_id)
        session_result = await self.db_session.execute(session_query)
        session = session_result.scalars().first()
        
        if not session:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Training session not found"
            )
        
        # Get coach info
        coach_query = select(User).where(User.id == session.coach_id)
        coach_result = await self.db_session.execute(coach_query)
        coach = coach_result.scalars().first()
        
        # Get feedbacks for this session
        feedback_query = select(Feedback).where(Feedback.training_session_id == session_id)
        feedback_result = await self.db_session.execute(feedback_query)
        feedbacks = feedback_result.scalars().all()
        
        # Calculate average ratings
        feedback_count = len(feedbacks)
        training_quality_avg = 0
        expectations_avg = 0
        body_condition_avg = 0
        intensity_avg = 0
        
        if feedback_count > 0:
            training_quality_avg = sum(fb.training_quality for fb in feedbacks) / feedback_count
            expectations_avg = sum(fb.expectations for fb in feedbacks) / feedback_count
            body_condition_avg = sum(fb.body_condition for fb in feedbacks) / feedback_count
            intensity_avg = sum(fb.intensity for fb in feedbacks) / feedback_count
        
        # Process each feedback
        feedback_details = []
        
        for feedback in feedbacks:
            # Get athlete info
            athlete_query = select(User).where(User.id == feedback.athlete_id)
            athlete_result = await self.db_session.execute(athlete_query)
            athlete = athlete_result.scalars().first()
            
            feedback_detail = FeedbackDetail(
                id=feedback.id,
                athlete_id=feedback.athlete_id,
                athlete_name=athlete.name if athlete else "Unknown",
                training_quality=feedback.training_quality,
                expectations=feedback.expectations,
                body_condition=feedback.body_condition,
                intensity=feedback.intensity,
                notes=feedback.notes,
                created_at=feedback.created_at.strftime("%Y-%m-%d %H:%M:%S")
            )
            
            feedback_details.append(feedback_detail)
        
        # Create session feedback report
        session_report = SessionFeedbackReport(
            session_id=session.id,
            date=session.date.strftime("%Y-%m-%d"),
            type=session.type,
            coach_name=coach.name if coach else "Unknown",
            feedback_count=feedback_count,
            training_quality_avg=round(training_quality_avg, 1),
            expectations_avg=round(expectations_avg, 1),
            body_condition_avg=round(body_condition_avg, 1),
            intensity_avg=round(intensity_avg, 1),
            feedbacks=feedback_details
        )
        
        # Create report summary
        summary = {
            "session_date": session.date.strftime("%Y-%m-%d"),
            "session_type": session.type,
            "coach_name": coach.name if coach else "Unknown",
            "feedback_count": feedback_count,
            "training_quality_avg": round(training_quality_avg, 1),
            "expectations_avg": round(expectations_avg, 1),
            "body_condition_avg": round(body_condition_avg, 1),
            "intensity_avg": round(intensity_avg, 1)
        }
        
        # Create the full report response
        report = FeedbackReportResponse(
            title=f"Feedback Report for {session.type} Session on {session.date.strftime('%Y-%m-%d')}",
            generated_at=datetime.now(),
            date_range=f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
            summary=summary,
            data=feedback_details
        )
        
        return report
    
    async def _generate_feedback_summary_report(
        self,
        start_date: datetime,
        end_date: datetime
    ) -> FeedbackReportResponse:
        """Generate a summary feedback report for all sessions in a date range"""
        # Get sessions in the date range
        session_query = select(TrainingSession).where(
            and_(
                TrainingSession.date >= start_date,
                TrainingSession.date <= end_date
            )
        ).order_by(TrainingSession.date.desc())
        
        session_result = await self.db_session.execute(session_query)
        sessions = session_result.scalars().all()
        
        # Process each session
        feedback_summaries = []
        total_feedback_count = 0
        total_training_quality = 0
        total_expectations = 0
        total_body_condition = 0
        total_intensity = 0
        
        for session in sessions:
            # Get coach info
            coach_query = select(User).where(User.id == session.coach_id)
            coach_result = await self.db_session.execute(coach_query)
            coach = coach_result.scalars().first()
            
            # Get feedbacks for this session
            feedback_query = select(Feedback).where(Feedback.training_session_id == session.id)
            feedback_result = await self.db_session.execute(feedback_query)
            feedbacks = feedback_result.scalars().all()
            
            # Calculate average ratings
            feedback_count = len(feedbacks)
            
            if feedback_count > 0:
                training_quality_avg = sum(fb.training_quality for fb in feedbacks) / feedback_count
                expectations_avg = sum(fb.expectations for fb in feedbacks) / feedback_count
                body_condition_avg = sum(fb.body_condition for fb in feedbacks) / feedback_count
                intensity_avg = sum(fb.intensity for fb in feedbacks) / feedback_count
                
                # Add to totals
                total_feedback_count += feedback_count
                total_training_quality += sum(fb.training_quality for fb in feedbacks)
                total_expectations += sum(fb.expectations for fb in feedbacks)
                total_body_condition += sum(fb.body_condition for fb in feedbacks)
                total_intensity += sum(fb.intensity for fb in feedbacks)
                
                # Create feedback summary
                feedback_summary = FeedbackSummary(
                    session_id=session.id,
                    date=session.date.strftime("%Y-%m-%d"),
                    type=session.type,
                    coach_name=coach.name if coach else "Unknown",
                    training_quality_avg=round(training_quality_avg, 1),
                    expectations_avg=round(expectations_avg, 1),
                    body_condition_avg=round(body_condition_avg, 1),
                    intensity_avg=round(intensity_avg, 1),
                    feedback_count=feedback_count
                )
                
                feedback_summaries.append(feedback_summary)
        
        # Calculate overall averages
        avg_training_quality = total_training_quality / total_feedback_count if total_feedback_count > 0 else 0
        avg_expectations = total_expectations / total_feedback_count if total_feedback_count > 0 else 0
        avg_body_condition = total_body_condition / total_feedback_count if total_feedback_count > 0 else 0
        avg_intensity = total_intensity / total_feedback_count if total_feedback_count > 0 else 0
        
        # Create report summary
        summary = {
            "total_sessions": len(sessions),
            "total_feedback": total_feedback_count,
            "avg_training_quality": round(avg_training_quality, 1),
            "avg_expectations": round(avg_expectations, 1),
            "avg_body_condition": round(avg_body_condition, 1),
            "avg_intensity": round(avg_intensity, 1)
        }
        
        # Create the full report response
        report = FeedbackReportResponse(
            title="Feedback Summary Report",
            generated_at=datetime.now(),
            date_range=f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
            summary=summary,
            data=feedback_summaries
        )
        
        return report
    
    async def generate_training_report_pdf(
        self,
        start_date: datetime,
        end_date: datetime,
        user_id: str,
        user_role: str
    ) -> bytes:
        """Generate a training report in PDF format"""
        # First get the report data
        report_data = await self.generate_training_report(
            start_date=start_date,
            end_date=end_date,
            user_id=user_id,
            user_role=user_role
        )
        
        # This is a placeholder for actual PDF generation
        # In a real implementation, you would use a library like ReportLab, WeasyPrint, or PDFKit
        # For now, we'll just create a simple text representation
        pdf_content = f"""
        3&7 Training & Recovery Platform
        {report_data.title}
        Generated: {report_data.generated_at.strftime('%Y-%m-%d %H:%M:%S')}
        Date Range: {report_data.date_range}
        
        Summary:
        Total Sessions: {report_data.summary['total_sessions']}
        Total Attendees: {report_data.summary['total_attendees']}
        Average Attendees Per Session: {report_data.summary['avg_attendees_per_session']}
        
        Sessions:
        """
        
        for session in report_data.data:
            pdf_content += f"""
            {session.date} - {session.type}
            Coach: {session.coach_name}
            Duration: {session.duration_minutes} minutes
            Attendees: {session.attendees_count}
            Feedback:
                Training Quality: {session.feedback['training_quality_avg']}/5
                Expectations: {session.feedback['expectations_avg']}/5
                Body Condition: {session.feedback['body_condition_avg']}/10
                Intensity: {session.feedback['intensity_avg']}/10
            """
        
        # Return the PDF content as bytes
        return pdf_content.encode('utf-8')
    
    async def generate_attendance_report_pdf(
        self,
        start_date: datetime,
        end_date: datetime,
        athlete_id: Optional[str] = None
    ) -> bytes:
        """Generate an attendance report in PDF format"""
        # First get the report data
        report_data = await self.generate_attendance_report(
            start_date=start_date,
            end_date=end_date,
            athlete_id=athlete_id
        )
        
        # Placeholder for PDF generation
        pdf_content = f"""
        3&7 Training & Recovery Platform
        {report_data.title}
        Generated: {report_data.generated_at.strftime('%Y-%m-%d %H:%M:%S')}
        Date Range: {report_data.date_range}
        
        Summary:
        """
        
        if "athlete_name" in report_data.summary:
            # Individual athlete report
            pdf_content += f"""
            Athlete: {report_data.summary['athlete_name']}
            Attended: {report_data.summary['attended_count']} of {report_data.summary['total_count']} sessions
            Attendance Rate: {report_data.summary['attendance_rate']}%
            
            Sessions:
            """
            
            for session in report_data.data[0].sessions:
                status = "Attended" if session.attended else "Missed"
                check_in = f" (Checked in at {session.check_in_time})" if session.attended and session.check_in_time else ""
                
                pdf_content += f"""
                {session.date} - {session.type}
                Coach: {session.coach_name}
                Status: {status}{check_in}
                """
        else:
            # Team report
            pdf_content += f"""
            Total Athletes: {report_data.summary['total_athletes']}
            Total Sessions: {report_data.summary['total_sessions']}
            Average Attendance Rate: {report_data.summary['avg_attendance_rate']}%
            
            Athletes:
            """
            
            for athlete in report_data.data:
                pdf_content += f"""
                {athlete.athlete_name}
                Attended: {athlete.sessions_attended} of {athlete.sessions_attended + athlete.sessions_missed} sessions
                Attendance Rate: {athlete.attendance_rate}%
                """
        
        # Return the PDF content as bytes
        return pdf_content.encode('utf-8')
    
    async def generate_feedback_report_pdf(
        self,
        start_date: datetime,
        end_date: datetime,
        session_id: Optional[str] = None
    ) -> bytes:
        """Generate a feedback report in PDF format"""
        # First get the report data
        report_data = await self.generate_feedback_report(
            start_date=start_date,
            end_date=end_date,
            session_id=session_id
        )
        
        # Placeholder for PDF generation
        pdf_content = f"""
        3&7 Training & Recovery Platform
        {report_data.title}
        Generated: {report_data.generated_at.strftime('%Y-%m-%d %H:%M:%S')}
        Date Range: {report_data.date_range}
        
        Summary:
        """
        
        if "session_date" in report_data.summary:
            # Session-specific report
            pdf_content += f"""
            Session Date: {report_data.summary['session_date']}
            Session Type: {report_data.summary['session_type']}
            Coach: {report_data.summary['coach_name']}
            Feedback Count: {report_data.summary['feedback_count']}
            
            Average Ratings:
            Training Quality: {report_data.summary['training_quality_avg']}/5
            Expectations: {report_data.summary['expectations_avg']}/5
            Body Condition: {report_data.summary['body_condition_avg']}/10
            Intensity: {report_data.summary['intensity_avg']}/10
            
            Individual Feedback:
            """
            
            for feedback in report_data.data:
                pdf_content += f"""
                Athlete: {feedback.athlete_name}
                Training Quality: {feedback.training_quality}/5
                Expectations: {feedback.expectations}/5
                Body Condition: {feedback.body_condition}/10
                Intensity: {feedback.intensity}/10
                Notes: {feedback.notes if feedback.notes else "None"}
                Submitted: {feedback.created_at}
                """
        else:
            # Summary report
            pdf_content += f"""
            Total Sessions: {report_data.summary['total_sessions']}
            Total Feedback: {report_data.summary['total_feedback']}
            
            Average Ratings:
            Training Quality: {report_data.summary['avg_training_quality']}/5
            Expectations: {report_data.summary['avg_expectations']}/5
            Body Condition: {report_data.summary['avg_body_condition']}/10
            Intensity: {report_data.summary['avg_intensity']}/10
            
            Sessions:
            """
            
            for session in report_data.data:
                pdf_content += f"""
                {session.date} - {session.type}
                Coach: {session.coach_name}
                Feedback Count: {session.feedback_count}
                Training Quality: {session.training_quality_avg}/5
                Expectations: {session.expectations_avg}/5
                Body Condition: {session.body_condition_avg}/10
                Intensity: {session.intensity_avg}/10
                """
        
        # Return the PDF content as bytes
        return pdf_content.encode('utf-8')
    
    async def generate_training_report_ppt(
        self,
        start_date: datetime,
        end_date: datetime,
        user_id: str,
        user_role: str
    ) -> bytes:
        """Generate a training report in PowerPoint format"""
        # First get the report data
        report_data = await self.generate_training_report(
            start_date=start_date,
            end_date=end_date,
            user_id=user_id,
            user_role=user_role
        )
        
        # Placeholder for PPT generation
        # In a real implementation, you would use a library like python-pptx
        # For now, we'll just return a message
        ppt_message = f"PowerPoint generation would be implemented here for {report_data.title}"
        
        # Return the message as bytes
        return ppt_message.encode('utf-8')
    
    async def generate_attendance_report_ppt(
        self,
        start_date: datetime,
        end_date: datetime,
        athlete_id: Optional[str] = None
    ) -> bytes:
        """Generate an attendance report in PowerPoint format"""
        # First get the report data
        report_data = await self.generate_attendance_report(
            start_date=start_date,
            end_date=end_date,
            athlete_id=athlete_id
        )
        
        # Placeholder for PPT generation
        ppt_message = f"PowerPoint generation would be implemented here for {report_data.title}"
        
        # Return the message as bytes
        return ppt_message.encode('utf-8')
    
    async def generate_feedback_report_ppt(
        self,
        start_date: datetime,
        end_date: datetime,
        session_id: Optional[str] = None
    ) -> bytes:
        """Generate a feedback report in PowerPoint format"""
        # First get the report data
        report_data = await self.generate_feedback_report(
            start_date=start_date,
            end_date=end_date,
            session_id=session_id
        )
        
        # Placeholder for PPT generation
        ppt_message = f"PowerPoint generation would be implemented here for {report_data.title}"
        
        # Return the message as bytes
        return ppt_message.encode('utf-8')
    
    async def export_custom_report_pdf(
        self,
        report_data: Dict[str, Any],
        title: str
    ) -> bytes:
        """Generate a custom report in PDF format based on provided data"""
        # Placeholder for PDF generation
        pdf_content = f"""
        3&7 Training & Recovery Platform
        {title}
        Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
        
        Custom Report Data:
        {json.dumps(report_data, indent=2)}
        """
        
        # Return the PDF content as bytes
        return pdf_content.encode('utf-8')
    
    async def export_custom_report_ppt(
        self,
        report_data: Dict[str, Any],
        title: str
    ) -> bytes:
        """Generate a custom report in PowerPoint format based on provided data"""
        # Placeholder for PPT generation
        ppt_message = f"PowerPoint generation would be implemented here for custom report: {title}"
        
        # Return the message as bytes
        return ppt_message.encode('utf-8')
    
    async def generate_monthly_report_pdf(
        self,
        start_date: datetime,
        end_date: datetime,
        user_id: str,
        user_role: str
    ) -> bytes:
        """Generate PDF for monthly report including independent training"""
        try:
            # Get training report data first
            training_report = await self.generate_training_report(start_date, end_date, user_id, user_role)
            
            # Fetch independent training data
            from repositories.independent_training_repository import IndependentTrainingRepository
            independent_repo = IndependentTrainingRepository(self.db_session)
            independent_training = await independent_repo.get_independent_training_sessions(
                start_date=start_date,
                end_date=end_date,
                user_id=user_id if user_role == 'athlete' else None,
                coach_id=user_id if user_role == 'coach' else None
            )
            
            # Create PDF content
            pdf_content = [
                # Title and introduction
                Paragraph(f"Monthly Training Report: {start_date.strftime('%B %Y')}", styles['Title']),
                Paragraph(f"Report generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}", styles['Normal']),
                Paragraph(f"Date range: {start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}", styles['Normal']),
                Spacer(1, 12),
                
                # Summary section
                Paragraph("Summary", styles['Heading1']),
                Paragraph(f"Total Training Sessions: {training_report.summary.get('total_sessions', 0)}", styles['Normal']),
                Paragraph(f"Total Attendees: {training_report.summary.get('total_attendees', 0)}", styles['Normal']),
                Paragraph(f"Average Attendees per Session: {training_report.summary.get('avg_attendees_per_session', 0):.1f}", styles['Normal']),
                Paragraph(f"Independent Training Sessions: {len(independent_training.sessions)}", styles['Normal']),
                Spacer(1, 12),
                
                # Group sessions section
                Paragraph("Group Training Sessions", styles['Heading1']),
            ]
            
            # Add group sessions table
            if training_report.data:
                table_data = [['Date', 'Type', 'Coach', 'Attendees', 'Duration (min)']]
                for session in training_report.data:
                    table_data.append([
                        session.date,
                        session.type,
                        session.coach_name,
                        str(session.attendees_count),
                        str(session.duration_minutes)
                    ])
                
                table = Table(table_data, colWidths=[80, 100, 100, 70, 80])
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                pdf_content.append(table)
            else:
                pdf_content.append(Paragraph("No group training sessions in this period", styles['Normal']))
            
            pdf_content.append(Spacer(1, 12))
            
            # Independent training section
            pdf_content.append(Paragraph("Independent Training Sessions", styles['Heading1']))
            
            if independent_training.sessions:
                ind_table_data = [['Date', 'Type', 'Location', 'Time', 'Intensity', 'Body Condition']]
                for session in independent_training.sessions:
                    ind_table_data.append([
                        session.date.strftime('%Y-%m-%d'),
                        session.type,
                        session.location,
                        f"{session.start_time.strftime('%H:%M')} - {session.end_time.strftime('%H:%M')}",
                        f"{session.intensity}/10",
                        f"{session.body_condition}/10"
                    ])
                
                ind_table = Table(ind_table_data, colWidths=[80, 80, 100, 80, 70, 70])
                ind_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                pdf_content.append(ind_table)
            else:
                pdf_content.append(Paragraph("No independent training sessions in this period", styles['Normal']))
            
            # Weekly breakdown section
            pdf_content.append(Spacer(1, 12))
            pdf_content.append(Paragraph("Weekly Breakdown", styles['Heading1']))
            
            # Process weeks - group by calendar week
            all_sessions = []
            
            # Add regular sessions with type
            for session in training_report.data:
                session_date = datetime.strptime(session.date, '%Y-%m-%d')
                week_num = session_date.isocalendar()[1]
                all_sessions.append({
                    'date': session_date,
                    'week': week_num,
                    'type': session.type,
                    'is_independent': False,
                    'title': f"Group: {session.type}",
                    'details': f"Coach: {session.coach_name}, Attendees: {session.attendees_count}"
                })
            
            # Add independent sessions with type
            for session in independent_training.sessions:
                week_num = session.date.isocalendar()[1]
                all_sessions.append({
                    'date': session.date,
                    'week': week_num,
                    'type': session.type,
                    'is_independent': True,
                    'title': f"Independent: {session.type}",
                    'details': f"Intensity: {session.intensity}/10, Body: {session.body_condition}/10"
                })
            
            # Group by week and create weekly tables
            weeks = {}
            for session in all_sessions:
                if session['week'] not in weeks:
                    weeks[session['week']] = []
                weeks[session['week']].append(session)
            
            # Create a table for each week
            for week_num, sessions in sorted(weeks.items()):
                pdf_content.append(Paragraph(f"Week {week_num}", styles['Heading2']))
                
                week_table_data = [['Date', 'Session Type', 'Details']]
                for session in sorted(sessions, key=lambda x: x['date']):
                    week_table_data.append([
                        session['date'].strftime('%Y-%m-%d'),
                        session['title'],
                        session['details']
                    ])
                
                week_table = Table(week_table_data, colWidths=[80, 150, 250])
                week_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                pdf_content.append(week_table)
                pdf_content.append(Spacer(1, 12))
            
            # Build the PDF and return
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            doc.build(pdf_content)
            
            return buffer.getvalue()
        except Exception as e:
            logger.error(f"Error generating monthly report PDF: {str(e)}")
            raise
    
    async def generate_monthly_report_ppt(
        self,
        start_date: datetime,
        end_date: datetime,
        user_id: str,
        user_role: str
    ) -> bytes:
        """Generate PPT for monthly report including independent training"""
        try:
            # Get training report data first
            training_report = await self.generate_training_report(start_date, end_date, user_id, user_role)
            
            # Fetch independent training data
            from repositories.independent_training_repository import IndependentTrainingRepository
            independent_repo = IndependentTrainingRepository(self.db_session)
            independent_training = await independent_repo.get_independent_training_sessions(
                start_date=start_date,
                end_date=end_date,
                user_id=user_id if user_role == 'athlete' else None,
                coach_id=user_id if user_role == 'coach' else None
            )
            
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"Monthly Training Report: {start_date.strftime('%B %Y')}"
            subtitle.text = f"Date range: {start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            
            # Summary slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "Summary"
            
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = "Training Overview"
            
            p = tf.add_paragraph()
            p.text = f"Group Sessions: {training_report.summary.get('total_sessions', 0)}"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = f"Total Attendees: {training_report.summary.get('total_attendees', 0)}"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = f"Average Attendees per Session: {training_report.summary.get('avg_attendees_per_session', 0):.1f}"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = f"Independent Training Sessions: {len(independent_training.sessions)}"
            p.level = 1
            
            # Group Sessions slide
            table_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(table_slide_layout)
            title = slide.shapes.title
            title.text = "Group Training Sessions"
            
            if training_report.data:
                rows = len(training_report.data) + 1  # +1 for header row
                cols = 5
                left = Inches(0.5)
                top = Inches(1.5)
                width = prs.slide_width - Inches(1)
                height = Inches(0.8 * min(rows, 8))  # Limit height for many rows
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column headings
                table.cell(0, 0).text = "Date"
                table.cell(0, 1).text = "Type"
                table.cell(0, 2).text = "Coach"
                table.cell(0, 3).text = "Attendees"
                table.cell(0, 4).text = "Duration (min)"
                
                # Apply heading formatting
                for cell in table.rows[0].cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                
                # Add data rows
                for idx, session in enumerate(training_report.data[:min(len(training_report.data), 7)]):  # Limit to 7 items
                    table.cell(idx+1, 0).text = session.date
                    table.cell(idx+1, 1).text = session.type
                    table.cell(idx+1, 2).text = session.coach_name
                    table.cell(idx+1, 3).text = str(session.attendees_count)
                    table.cell(idx+1, 4).text = str(session.duration_minutes)
            else:
                left = Inches(1)
                top = Inches(2)
                width = prs.slide_width - Inches(2)
                height = Inches(1)
                
                txbox = slide.shapes.add_textbox(left, top, width, height)
                tf = txbox.text_frame
                tf.text = "No group training sessions in this period"
            
            # Independent Training slide
            slide = prs.slides.add_slide(table_slide_layout)
            title = slide.shapes.title
            title.text = "Independent Training Sessions"
            
            if independent_training.sessions:
                rows = min(len(independent_training.sessions) + 1, 8)  # +1 for header, max 8 rows
                cols = 6
                left = Inches(0.5)
                top = Inches(1.5)
                width = prs.slide_width - Inches(1)
                height = Inches(0.8 * (rows - 1))
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column headings
                table.cell(0, 0).text = "Date"
                table.cell(0, 1).text = "Type"
                table.cell(0, 2).text = "Location"
                table.cell(0, 3).text = "Time"
                table.cell(0, 4).text = "Intensity"
                table.cell(0, 5).text = "Body Cond."
                
                # Apply heading formatting
                for cell in table.rows[0].cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                
                # Add data rows
                for idx, session in enumerate(independent_training.sessions[:rows-1]):
                    table.cell(idx+1, 0).text = session.date.strftime('%Y-%m-%d')
                    table.cell(idx+1, 1).text = session.type
                    table.cell(idx+1, 2).text = session.location
                    table.cell(idx+1, 3).text = f"{session.start_time.strftime('%H:%M')} - {session.end_time.strftime('%H:%M')}"
                    table.cell(idx+1, 4).text = f"{session.intensity}/10"
                    table.cell(idx+1, 5).text = f"{session.body_condition}/10"
            else:
                left = Inches(1)
                top = Inches(2)
                width = prs.slide_width - Inches(2)
                height = Inches(1)
                
                txbox = slide.shapes.add_textbox(left, top, width, height)
                tf = txbox.text_frame
                tf.text = "No independent training sessions in this period"
            
            # Weekly Breakdown slides
            # Process weeks - group by calendar week
            all_sessions = []
            
            # Add regular sessions with type
            for session in training_report.data:
                session_date = datetime.strptime(session.date, '%Y-%m-%d')
                week_num = session_date.isocalendar()[1]
                all_sessions.append({
                    'date': session_date,
                    'week': week_num,
                    'type': session.type,
                    'is_independent': False,
                    'title': f"Group: {session.type}",
                    'details': f"Coach: {session.coach_name}, Attendees: {session.attendees_count}"
                })
            
            # Add independent sessions with type
            for session in independent_training.sessions:
                week_num = session.date.isocalendar()[1]
                all_sessions.append({
                    'date': session.date,
                    'week': week_num,
                    'type': session.type,
                    'is_independent': True,
                    'title': f"Independent: {session.type}",
                    'details': f"Intensity: {session.intensity}/10, Body: {session.body_condition}/10"
                })
            
            # Group by week
            weeks = {}
            for session in all_sessions:
                if session['week'] not in weeks:
                    weeks[session['week']] = []
                weeks[session['week']].append(session)
            
            # Create slides for each week (if not too many)
            for week_num in sorted(weeks.keys())[:4]:  # Limit to first 4 weeks
                sessions = weeks[week_num]
                
                slide = prs.slides.add_slide(table_slide_layout)
                title = slide.shapes.title
                title.text = f"Week {week_num} Breakdown"
                
                rows = min(len(sessions) + 1, 10)  # +1 for header, max 10 rows
                cols = 3
                left = Inches(0.5)
                top = Inches(1.5)
                width = prs.slide_width - Inches(1)
                height = Inches(0.5 * (rows))
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column headings
                table.cell(0, 0).text = "Date"
                table.cell(0, 1).text = "Session Type"
                table.cell(0, 2).text = "Details"
                
                # Apply heading formatting
                for cell in table.rows[0].cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                
                # Add data rows
                sorted_sessions = sorted(sessions, key=lambda x: x['date'])
                for idx, session in enumerate(sorted_sessions[:rows-1]):
                    table.cell(idx+1, 0).text = session['date'].strftime('%Y-%m-%d')
                    table.cell(idx+1, 1).text = session['title']
                    table.cell(idx+1, 2).text = session['details']
                    
                    # Highlight independent sessions
                    if session['is_independent']:
                        table.cell(idx+1, 1).fill.solid()
                        table.cell(idx+1, 1).fill.fore_color.rgb = RGBColor(242, 220, 179)  # Light orange
            
            # Save the presentation
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return buffer.getvalue()
        except Exception as e:
            logger.error(f"Error generating monthly report PPT: {str(e)}")
            raise 